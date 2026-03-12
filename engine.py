"""
XlArch Engine — One architecture dict → canvas, PPTX, PNG.
All renderers share the same coordinate system.
"""

import os, base64, io, json, subprocess, tempfile
from pathlib import Path

ICON_DIR = Path(__file__).parent / "static" / "icons"

def get_icon_path(key):
    p = ICON_DIR / f"{key}.png"
    return str(p) if p.exists() else None

def get_all_icon_keys():
    return sorted([p.stem for p in ICON_DIR.glob("*.png")])

CATALOG = {
    "GCP Compute": [
        {"id": "cloud_run", "label": "Cloud Run", "icon": "cloud_run"},
        {"id": "functions", "label": "Cloud Functions", "icon": "functions"},
        {"id": "gke", "label": "GKE", "icon": "gke"},
        {"id": "compute_engine", "label": "Compute Engine", "icon": "compute_engine"},
    ],
    "GCP Analytics": [
        {"id": "bigquery", "label": "BigQuery", "icon": "bigquery"},
        {"id": "dataflow", "label": "Dataflow", "icon": "dataflow"},
        {"id": "pubsub", "label": "Pub/Sub", "icon": "pubsub"},
        {"id": "composer", "label": "Composer", "icon": "composer"},
        {"id": "looker", "label": "Looker", "icon": "looker"},
        {"id": "dataproc", "label": "Dataproc", "icon": "dataproc"},
        {"id": "data_catalog", "label": "Data Catalog", "icon": "data_catalog"},
    ],
    "GCP Storage & DB": [
        {"id": "gcs", "label": "Cloud Storage", "icon": "gcs"},
        {"id": "bigtable", "label": "Bigtable", "icon": "bigtable"},
        {"id": "memorystore", "label": "Memorystore", "icon": "memorystore"},
        {"id": "spanner", "label": "Spanner", "icon": "spanner"},
        {"id": "firestore", "label": "Firestore", "icon": "firestore"},
        {"id": "cloudsql", "label": "Cloud SQL", "icon": "cloudsql"},
    ],
    "GCP AI / ML": [
        {"id": "vertex_ai", "label": "Vertex AI", "icon": "vertex_ai"},
        {"id": "automl", "label": "AutoML", "icon": "automl"},
    ],
    "GCP Security & Ops": [
        {"id": "iam", "label": "Cloud IAM", "icon": "iam"},
        {"id": "kms", "label": "Cloud KMS", "icon": "kms"},
        {"id": "monitoring", "label": "Monitoring", "icon": "monitoring"},
        {"id": "logging", "label": "Logging", "icon": "logging"},
    ],
    "GCP Network": [
        {"id": "load_balancing", "label": "Load Balancing", "icon": "load_balancing"},
        {"id": "cdn", "label": "Cloud CDN", "icon": "cdn"},
        {"id": "dns", "label": "Cloud DNS", "icon": "dns"},
    ],
    "Databases": [
        {"id": "postgresql", "label": "PostgreSQL", "icon": "postgresql"},
        {"id": "mysql", "label": "MySQL", "icon": "mysql"},
        {"id": "oracle", "label": "Oracle", "icon": "oracle"},
        {"id": "mssql", "label": "SQL Server", "icon": "mssql"},
        {"id": "mongodb", "label": "MongoDB", "icon": "mongodb"},
        {"id": "cassandra", "label": "Cassandra", "icon": "cassandra"},
        {"id": "redis", "label": "Redis", "icon": "redis"},
        {"id": "generic_db", "label": "Database", "icon": "generic_db"},
    ],
    "Messaging": [
        {"id": "kafka", "label": "Kafka", "icon": "kafka"},
        {"id": "rabbitmq", "label": "RabbitMQ", "icon": "rabbitmq"},
        {"id": "nats", "label": "NATS", "icon": "nats"},
    ],
    "AWS": [
        {"id": "aurora", "label": "Aurora", "icon": "aurora"},
        {"id": "dynamodb", "label": "DynamoDB", "icon": "dynamodb"},
        {"id": "redshift", "label": "Redshift", "icon": "redshift"},
        {"id": "rds", "label": "RDS", "icon": "rds"},
    ],
    "General": [
        {"id": "user", "label": "User", "icon": "user"},
        {"id": "users", "label": "Users", "icon": "users"},
        {"id": "client", "label": "Client App", "icon": "client"},
        {"id": "database", "label": "Database", "icon": "database"},
    ],
}

# ═══ SHARED LAYOUT CONSTANTS ═══
DEFAULT_ZONES = ["sources", "ingest", "process", "store", "analyze", "serve"]
DEFAULT_LANES = ["streaming", "batch", "ml"]

NW = 120     # node width
NH = 82      # node height
IS = 42      # icon size
ZG = 175     # zone gap
LG = 145     # lane gap
LP = 70      # left pad
TP = 90      # top pad
GOV_OFF = 50 # governance offset below last lane
ARROW_COLOR = "#9E9E9E"
LANE_COLORS = {"streaming": "#E8F4FD", "batch": "#EDF7ED", "ml": "#FFF3E0"}


def compute_layout(architecture):
    zones = architecture.get("zones", DEFAULT_ZONES)
    lanes = architecture.get("lanes", DEFAULT_LANES)
    nodes = architecture.get("nodes", [])

    zone_x = {z: LP + i * ZG for i, z in enumerate(zones)}
    lane_y = {l: TP + i * LG for i, l in enumerate(lanes)}

    cells = {}
    for n in nodes:
        key = f"{n.get('zone','')}-{n.get('lane','')}"
        cells.setdefault(key, []).append(n["id"])

    positions = {}
    for n in nodes:
        z, l = n.get("zone", ""), n.get("lane", "")
        key = f"{z}-{l}"
        siblings = cells.get(key, [n["id"]])
        idx = siblings.index(n["id"]) if n["id"] in siblings else 0
        positions[n["id"]] = {
            "x": zone_x.get(z, LP) + idx * (NW + 30),
            "y": lane_y.get(l, TP),
        }
    return positions


def canvas_bounds(architecture):
    zones = architecture.get("zones", DEFAULT_ZONES)
    lanes = architecture.get("lanes", DEFAULT_LANES)
    return LP + len(zones) * ZG + 60, TP + len(lanes) * LG + GOV_OFF + 100


def new_architecture(title="Untitled"):
    return {
        "title": title,
        "zones": list(DEFAULT_ZONES),
        "lanes": list(DEFAULT_LANES),
        "nodes": [], "edges": [], "governance": [], "bestPractices": [],
    }


# ═══ SYSTEM PROMPT ═══
SYSTEM_PROMPT = """You are XlArch, an expert system architect. Generate architecture as a Python dict.

OUTPUT ONLY A VALID PYTHON DICT LITERAL. No markdown, no backticks, no explanation.
Start with { and end with }

Available icon IDs: """ + ", ".join(get_all_icon_keys()) + """

ZONES (left-to-right): sources, ingest, process, store, analyze, serve
LANES (top-to-bottom): streaming, batch, ml

FORMAT:
{
    "title": "Title",
    "zones": ["sources", "ingest", "process", "store", "analyze", "serve"],
    "lanes": ["streaming", "batch", "ml"],
    "nodes": [
        {"id": "unique_id", "icon": "bigquery", "label": "BigQuery", "zone": "store", "lane": "batch", "step": 5},
    ],
    "edges": [
        {"from": "id1", "to": "id2"},
    ],
    "governance": [
        {"icon": "iam", "label": "Cloud IAM"},
        {"icon": "kms", "label": "Encryption (CMEK)"},
        {"icon": "monitoring", "label": "Monitoring"},
        {"icon": "logging", "label": "Logging & Audit"},
        {"icon": "firewall", "label": "VPC-SC / Blast Radius"},
        {"icon": "dns", "label": "Private Connectivity"},
    ],
    "bestPractices": [
        {"category": "SECURITY", "tip": "..."},
    ],
}

═══════════════════════════════════════
NON-NEGOTIABLE — EVERY architecture MUST include these in governance:
═══════════════════════════════════════

1. IDENTITY & ACCESS
   - Cloud IAM (icon: iam) — least privilege, service accounts, Workload Identity
   - Authentication & Authorization flow must be shown or referenced

2. ENCRYPTION
   - Cloud KMS (icon: kms) — CMEK for all data at rest
   - TLS/mTLS for all data in transit

3. NETWORK SECURITY
   - VPC Service Controls / blast radius isolation (icon: firewall)
   - Private connectivity — Private Service Connect, VPC peering, no public endpoints (icon: dns)
   - Interconnect for hybrid connectivity if applicable

4. OBSERVABILITY
   - Cloud Monitoring (icon: monitoring) — metrics, dashboards, SLOs
   - Cloud Logging (icon: logging) — audit logs, access transparency, log sinks
   - Alerting policies on all critical services

5. DATA GOVERNANCE
   - Data Catalog / Dataplex (icon: data_catalog) — metadata, lineage, classification
   - DLP for PII detection where applicable

ALL of the above MUST appear in the "governance" array. Never skip them.

═══════════════════════════════════════
ARCHITECTURE RULES:
═══════════════════════════════════════

- 10-18 nodes in the main architecture
- Number steps 1-N in data flow order
- Every edge represents a real data flow — show HOW data moves
- Use only icon IDs from the available list
- Each node in exactly one zone and one lane
- Connect logically: sources → ingest → process → store → analyze → serve
- Show cross-lane connections where data flows between streaming/batch/ml paths

BEST PRACTICES — include 6-10 across these categories:
- SECURITY: IAM least privilege, CMEK, VPC-SC, no public endpoints, service account keys rotation
- NETWORK: Private Service Connect, dedicated interconnect, firewall rules, DNS peering
- RELIABILITY: multi-region, auto-scaling, dead letter queues, retry policies, HA configurations
- COST: committed use discounts, autoscaling to zero, lifecycle policies, slot reservations
- PERFORMANCE: caching, partitioning, pre-splitting, streaming engine, connection pooling
- COMPLIANCE: audit logging, access transparency, data residency, retention policies
"""


# ═══ PPTX RENDERER ═══

def render_pptx(architecture, output_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    positions = compute_layout(architecture)
    cw, ch = canvas_bounds(architecture)
    zones = architecture.get("zones", DEFAULT_ZONES)
    lanes = architecture.get("lanes", DEFAULT_LANES)

    # Scale virtual canvas → 13.3" × 7.5" slide
    margin = 0.35
    scale = min((13.3 - 2*margin) / cw, (7.5 - 2*margin) / ch)
    def sx(x): return margin + x * scale
    def sy(y): return margin + y * scale
    def sw(v): return v * scale

    prs = Presentation()
    prs.slide_width = Inches(13.3)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:],16))

    def txt(x, y, w, h, text, sz=7, col="3C4043", bold=False, align="center"):
        tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tx.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        p.text = text; p.font.size = Pt(sz); p.font.color.rgb = rgb(col); p.font.bold = bold
        p.alignment = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT}.get(align, PP_ALIGN.CENTER)

    # Title
    txt(sx(0), 0.12, sw(cw), 0.35, architecture.get("title", ""), sz=14, col="202124", bold=True, align="left")

    # Lane backgrounds
    for i, lane in enumerate(lanes):
        ly = TP + i * LG - 18
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(sx(25)), Inches(sy(ly)), Inches(sw(len(zones)*ZG+40)), Inches(sw(NH+50)))
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(LANE_COLORS.get(lane, "#F5F5F5"))
        shape.line.fill.background()

    # Zone labels
    for i, z in enumerate(zones):
        txt(sx(LP + i*ZG), sy(25), sw(NW), 0.15, z.upper(), sz=6.5, col="BDBDBD", bold=True)

    # Nodes
    nw_in, nh_in, is_in = sw(NW), sw(NH), sw(IS)
    for n in architecture.get("nodes", []):
        p = positions.get(n["id"])
        if not p: continue
        px, py = sx(p["x"]), sy(p["y"])

        # Icon
        ipath = get_icon_path(n.get("icon", ""))
        if ipath:
            try: slide.shapes.add_picture(ipath, Inches(px + (nw_in-is_in)/2), Inches(py+0.02), Inches(is_in), Inches(is_in))
            except: pass

        # Label
        txt(px-0.03, py+is_in+0.01, nw_in+0.06, 0.22, n.get("label",""), sz=6.5)

        # Step badge
        if "step" in n:
            bsz = 0.17
            s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(px+nw_in-0.01), Inches(py-0.02), Inches(bsz), Inches(bsz))
            s.fill.solid(); s.fill.fore_color.rgb = rgb("4285F4"); s.line.fill.background()
            p2 = s.text_frame.paragraphs[0]
            p2.text = str(n["step"]); p2.font.size = Pt(7)
            p2.font.color.rgb = RGBColor(255,255,255); p2.font.bold = True
            p2.alignment = PP_ALIGN.CENTER

    # Edges — orthogonal routing
    for e in architecture.get("edges", []):
        fp, tp = positions.get(e["from"]), positions.get(e["to"])
        if not fp or not tp: continue
        x1, y1 = fp["x"]+NW, fp["y"]+NH/2
        x2, y2 = tp["x"], tp["y"]+NH/2

        if abs(y1 - y2) < 5:
            # Same lane — straight
            conn = slide.shapes.add_connector(1,
                Inches(sx(x1)), Inches(sy(y1)), Inches(sx(x2)), Inches(sy(y2)))
            conn.line.color.rgb = rgb(ARROW_COLOR); conn.line.width = Pt(0.75)
        else:
            # Cross-lane — draw 3 segments (H, V, H)
            midX = (x1 + x2) / 2
            for (ax, ay, bx, by) in [(x1,y1,midX,y1), (midX,y1,midX,y2), (midX,y2,x2,y2)]:
                c = slide.shapes.add_connector(1,
                    Inches(sx(ax)), Inches(sy(ay)), Inches(sx(bx)), Inches(sy(by)))
                c.line.color.rgb = rgb(ARROW_COLOR); c.line.width = Pt(0.75)

    # Governance bar
    gov = architecture.get("governance", [])
    if gov:
        gy = TP + len(lanes)*LG + GOV_OFF
        bw = len(zones)*ZG + 40
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(sx(25)), Inches(sy(gy-8)), Inches(sw(bw)), Inches(sw(80)))
        s.fill.solid(); s.fill.fore_color.rgb = rgb("F3E8FD")
        s.line.color.rgb = rgb("A142F4"); s.line.width = Pt(0.3)

        txt(sx(35), sy(gy-3), sw(300), 0.12, "GOVERNANCE · SECURITY · OBSERVABILITY", sz=5, col="8430CE", bold=True, align="left")

        gis = sw(24)
        usable = bw - 60
        spacing = usable / max(len(gov), 1)
        for i, g in enumerate(gov):
            gx = 45 + i * spacing
            ipath = get_icon_path(g.get("icon",""))
            if ipath:
                try: slide.shapes.add_picture(ipath, Inches(sx(gx)), Inches(sy(gy+12)), Inches(gis), Inches(gis))
                except: pass
            label = g.get("label","")
            if len(label) > 16: label = label[:14] + "…"
            txt(sx(gx)-0.06, sy(gy+12)+gis+0.01, gis+0.12, 0.12, label, sz=4.5, col="5F6368")

    prs.save(output_path)
    return output_path


# ═══ PNG RENDERER ═══

def render_png(architecture, output_path):
    nodes = architecture.get("nodes", [])
    edges = architecture.get("edges", [])
    title = architecture.get("title", "Architecture")

    DMAP = {
        "bigquery":("diagrams.gcp.analytics","BigQuery"), "dataflow":("diagrams.gcp.analytics","Dataflow"),
        "pubsub":("diagrams.gcp.analytics","PubSub"), "composer":("diagrams.gcp.analytics","Composer"),
        "looker":("diagrams.gcp.analytics","Looker"), "dataproc":("diagrams.gcp.analytics","Dataproc"),
        "gcs":("diagrams.gcp.storage","GCS"), "bigtable":("diagrams.gcp.database","BigTable"),
        "memorystore":("diagrams.gcp.database","Memorystore"), "spanner":("diagrams.gcp.database","Spanner"),
        "firestore":("diagrams.gcp.database","Firestore"), "cloudsql":("diagrams.gcp.database","SQL"),
        "vertex_ai":("diagrams.gcp.ml","VertexAI"), "cloud_run":("diagrams.gcp.compute","CloudRun"),
        "functions":("diagrams.gcp.compute","Functions"), "gke":("diagrams.gcp.compute","GKE"),
        "compute_engine":("diagrams.gcp.compute","ComputeEngine"),
        "iam":("diagrams.gcp.security","Iam"), "kms":("diagrams.gcp.security","KMS"),
        "monitoring":("diagrams.gcp.operations","Monitoring"), "logging":("diagrams.gcp.operations","Logging"),
        "kafka":("diagrams.onprem.queue","Kafka"), "rabbitmq":("diagrams.onprem.queue","RabbitMQ"),
        "postgresql":("diagrams.onprem.database","PostgreSQL"), "mysql":("diagrams.onprem.database","MySQL"),
        "oracle":("diagrams.onprem.database","Oracle"), "mssql":("diagrams.onprem.database","Mssql"),
        "mongodb":("diagrams.onprem.database","MongoDB"), "redis":("diagrams.onprem.inmemory","Redis"),
    }

    imports, nmap = set(), {}
    for n in nodes:
        if n.get("icon","") in DMAP:
            mod, cls = DMAP[n["icon"]]
            imports.add(f"from {mod} import {cls}")
            nmap[n["id"]] = (cls, n.get("label", n["id"]))

    if not nmap: return None

    lines = ["from diagrams import Diagram, Edge", *sorted(imports), "",
        f'with Diagram("{title}", filename="{output_path}", show=False, direction="LR",',
        '             graph_attr={"splines":"ortho","nodesep":"0.6","ranksep":"0.9","bgcolor":"#FFFFFF"}):']
    for nid, (cls, label) in nmap.items():
        s = nid.replace("-","_").replace(" ","_")
        lines.append(f'    {s} = {cls}("{label}")')
    valid = {nid.replace("-","_").replace(" ","_") for nid in nmap}
    for e in edges:
        f, t = e["from"].replace("-","_").replace(" ","_"), e["to"].replace("-","_").replace(" ","_")
        if f in valid and t in valid: lines.append(f'    {f} >> {t}')

    with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False) as f:
        f.write("\n".join(lines)); sp = f.name
    try:
        r = subprocess.run(["python3", sp], capture_output=True, text=True, timeout=30)
        if r.returncode == 0: return f"{output_path}.png"
    except: pass
    finally: os.unlink(sp)
    return None
