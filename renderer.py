"""
XlArch Renderer — ONE render function, multiple output adapters.
render() produces draw commands. Adapters translate to canvas JS or PPTX shapes.
"""

from engine import compute_layout, get_icon_path, get_all_icon_keys, NW, NH, IS, ZG, LG, LP, TP, GOV_OFF, LANE_COLORS, ARROW_COLOR

# ═══════════════════════════════════════════
#  DRAW COMMANDS — the universal language
# ═══════════════════════════════════════════
# Each command is a dict: {"type": "...", ...params}
# Types:
#   rect:  {x, y, w, h, fill, stroke, stroke_width, radius}
#   image: {x, y, w, h, src (icon key)}
#   text:  {x, y, w, text, size, color, bold, align}
#   line:  {x1, y1, x2, y2, color, width}
#   arrow: {points: [(x,y)...], color, width}  — polyline with arrowhead at end


def render(architecture):
    """Render architecture to a list of draw commands. Canvas + PPTX both consume this."""
    commands = []

    arch = dict(architecture)
    zones = list(arch.get("zones", []))
    lanes = list(arch.get("lanes", []))
    nodes = list(arch.get("nodes", []))
    edges = list(arch.get("edges", []))

    arch["zones"] = zones
    arch["lanes"] = lanes
    arch["nodes"] = nodes
    arch["edges"] = edges

    positions, zone_gap, lane_gap = compute_layout(arch)
    title = arch.get("title", "")

    # Canvas size from actual node positions (handles overflow)
    max_x = max((p["x"] for p in positions.values()), default=0) + NW + 60
    cw = max(LP + len(zones) * zone_gap + 60, max_x)
    ch = TP + len(lanes) * lane_gap + 60

    # ── Title
    commands.append({"type": "text", "x": LP, "y": 15, "w": cw, "text": title,
                     "size": 14, "color": "#202124", "bold": True, "align": "left"})

    # ── Lane backgrounds + lane labels
    for i, lane in enumerate(lanes):
        ly = TP + i * lane_gap - 20
        lh = NH + 58
        
        # Calculate lane width from actual node positions
        lane_nodes = [n for n in nodes if n.get("lane") == lane]
        if lane_nodes:
            max_x = max(positions.get(n["id"], {}).get("x", 0) for n in lane_nodes)
            lw = max(len(zones) * zone_gap + 40, max_x + NW + 40 - 25)
        else:
            lw = len(zones) * zone_gap + 40
        
        fill = LANE_COLORS.get(lane, "#F5F5F5")
        commands.append({"type": "rect", "x": 25, "y": ly, "w": lw, "h": lh,
                        "fill": fill, "stroke": "#E8E8E8", "stroke_width": 0.3, "radius": 10, "opacity": 0.35})

        # Lane label
        label = lane.upper().replace("_", " ")
        commands.append({"type": "text", "x": 30, "y": ly + lh / 2 - 5, "w": 60,
                        "text": label, "size": 6, "color": "#BDBDBD", "bold": True, "align": "left"})

    # ── Zone labels
    for i, z in enumerate(zones):
        label = z.upper().replace("_", " ")
        commands.append({"type": "text", "x": LP + i * zone_gap, "y": 50, "w": NW + 20,
                        "text": label, "size": 7, "color": "#9E9E9E", "bold": True, "align": "center"})

    # ── Edges (smart orthogonal routing)
    # Track used midpoints to offset parallel arrows
    used_mid = {}

    for e in edges:
        fp, tp = positions.get(e["from"]), positions.get(e["to"])
        if not fp or not tp:
            continue

        fx, fy = fp["x"], fp["y"]
        tx, ty = tp["x"], tp["y"]
        fcx, fcy = fx + NW / 2, fy + NH / 2  # from center
        tcx, tcy = tx + NW / 2, ty + NH / 2  # to center

        dx = tcx - fcx
        dy = tcy - fcy

        GAP = 8  # gap from node edge

        if abs(dy) < 10:
            # Same lane — horizontal
            if dx > 0:
                # Left to right
                points = [(fx + NW + GAP, fcy), (tx - GAP, tcy)]
            else:
                # Right to left — route around bottom
                bot = max(fy, ty) + NH + 25
                points = [(fx + NW + GAP, fcy), (fx + NW + GAP, bot),
                           (tx - GAP, bot), (tx - GAP, tcy)]

        elif abs(dx) < 10:
            # Same zone — vertical
            if dy > 0:
                points = [(fcx, fy + NH + GAP), (tcx, ty - GAP)]
            else:
                points = [(fcx, fy - GAP), (tcx, ty + NH + GAP)]

        elif dx > 0:
            # Forward + cross-lane — L-shape exit right, enter left
            mid_key = round((fx + NW + tx) / 2 / 20) * 20
            offset = used_mid.get(mid_key, 0) * 12
            used_mid[mid_key] = used_mid.get(mid_key, 0) + 1
            midX = (fx + NW + tx) / 2 + offset

            points = [(fx + NW + GAP, fcy), (midX, fcy), (midX, tcy), (tx - GAP, tcy)]

        else:
            # Backward + cross-lane — route around
            if dy > 0:
                midY = fy + NH + 25
                points = [(fx + NW + GAP, fcy), (fx + NW + GAP, midY),
                           (tx - GAP, midY), (tx - GAP, tcy)]
            else:
                midY = ty + NH + 25
                points = [(fcx, fy - GAP), (fcx, midY),
                           (tx - GAP, midY), (tx - GAP, tcy)]

        commands.append({"type": "arrow", "points": points,
                        "color": ARROW_COLOR, "width": 2,
                        "fromId": e["from"], "toId": e["to"]})

    # ── Nodes
    for n in nodes:
        p = positions.get(n["id"])
        if not p:
            continue
        x, y = p["x"], p["y"]
        nid = n["id"]

        # Card background
        commands.append({"type": "rect", "x": x, "y": y, "w": NW, "h": NH,
                        "fill": "#FFFFFF", "stroke": "#E0E4E8", "stroke_width": 1,
                        "radius": 10, "shadow": True, "nodeId": nid})

        # Icon (bigger, centered)
        icon_size = 54
        commands.append({"type": "image", "x": x + (NW - icon_size) / 2, "y": y + 6,
                        "w": icon_size, "h": icon_size, "src": n.get("icon", ""), "nodeId": nid})

        # Label
        label = n.get("label", "")
        if len(label) > 18:
            label = label[:16] + "…"
        commands.append({"type": "text", "x": x, "y": y + NH - 16, "w": NW,
                        "text": label, "size": 9, "color": "#3C4043",
                        "bold": False, "align": "center", "nodeId": nid})

        # Step badge
        if "step" in n:
            bx, by, bsz = x + NW - 12, y - 6, 18
            commands.append({"type": "rect", "x": bx, "y": by, "w": bsz, "h": 16,
                            "fill": "#4285F4", "stroke": None, "stroke_width": 0,
                            "radius": 8, "nodeId": nid})
            commands.append({"type": "text", "x": bx, "y": by + 1, "w": bsz,
                            "text": str(n["step"]), "size": 8, "color": "#FFFFFF",
                            "bold": True, "align": "center", "nodeId": nid})

    return {"commands": commands, "width": cw, "height": ch, "positions": positions, "architecture": arch}


# ═══════════════════════════════════════════
#  ADAPTER: PPTX
# ═══════════════════════════════════════════

def commands_to_pptx(render_result, output_path):
    """Translate draw commands to python-pptx shapes, grouped by nodeId."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    commands = render_result["commands"]
    cw, ch = render_result["width"], render_result["height"]

    SLIDE_W, SLIDE_H = 13.3, 7.5
    margin = 0.35
    scale = min((SLIDE_W - 2 * margin) / cw, (SLIDE_H - 2 * margin) / ch)

    def sx(x): return margin + x * scale
    def sy(y): return margin + y * scale
    def sw(v): return v * scale

    def rgb(h):
        h = h.lstrip("#")
        if len(h) == 3:
            h = h[0]*2 + h[1]*2 + h[2]*2
        return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Collect shapes per nodeId for grouping
    node_shapes = {}  # nodeId → [shape, ...]

    def add_shape_for_cmd(cmd):
        """Create a shape and return it."""
        t = cmd["type"]
        shape = None

        if t == "rect":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE if cmd.get("radius", 0) > 0 else MSO_SHAPE.RECTANGLE,
                Inches(sx(cmd["x"])), Inches(sy(cmd["y"])),
                Inches(sw(cmd["w"])), Inches(sw(cmd["h"])))
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(cmd["fill"])
            if cmd.get("stroke"):
                shape.line.color.rgb = rgb(cmd["stroke"])
                shape.line.width = Pt(cmd.get("stroke_width", 0.5))
            else:
                shape.line.fill.background()
            if cmd.get("shadow") is not True:
                shape.shadow.inherit = False

        elif t == "image":
            ipath = get_icon_path(cmd["src"])
            if ipath:
                try:
                    shape = slide.shapes.add_picture(ipath,
                        Inches(sx(cmd["x"])), Inches(sy(cmd["y"])),
                        Inches(sw(cmd["w"])), Inches(sw(cmd["h"])))
                except:
                    pass

        elif t == "text":
            shape = slide.shapes.add_textbox(
                Inches(sx(cmd["x"])), Inches(sy(cmd["y"])),
                Inches(sw(cmd["w"])), Inches(0.25))
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cmd["text"]
            p.font.size = Pt(cmd.get("size", 7))
            p.font.color.rgb = rgb(cmd.get("color", "#3C4043"))
            p.font.bold = cmd.get("bold", False)
            align_map = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}
            p.alignment = align_map.get(cmd.get("align", "center"), PP_ALIGN.CENTER)

        elif t == "arrow":
            pts = cmd["points"]
            color = rgb(cmd.get("color", "#9E9E9E"))
            width = Pt(cmd.get("width", 2))
            for j in range(len(pts) - 1):
                x1, y1 = pts[j]
                x2, y2 = pts[j + 1]
                conn = slide.shapes.add_connector(1,
                    Inches(sx(x1)), Inches(sy(y1)),
                    Inches(sx(x2)), Inches(sy(y2)))
                conn.line.color.rgb = color
                conn.line.width = width

        return shape

    # First pass: create all shapes
    for cmd in commands:
        shape = add_shape_for_cmd(cmd)
        nid = cmd.get("nodeId")
        if nid and shape:
            node_shapes.setdefault(nid, []).append(shape)

    # Second pass: group shapes by nodeId
    for nid, shapes in node_shapes.items():
        if len(shapes) >= 2:
            try:
                group = slide.shapes.add_group_shape(shapes)
            except:
                pass  # grouping may fail on some shape combos, shapes stay ungrouped

    prs.save(output_path)
    return output_path


# ═══════════════════════════════════════════
#  ADAPTER: CANVAS JS
# ═══════════════════════════════════════════

def commands_to_canvas_js(render_result):
    """Translate draw commands to JavaScript that draws on an HTML canvas."""
    import json
    return json.dumps(render_result["commands"])
